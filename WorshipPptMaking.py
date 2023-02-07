import os
import shutil
import pandas as pd
import datetime


# update


# 예배에 쓸 찬송가 복사 함수
def song_copy(worship_songs_copy):
    path = "./새찬송가PPT/"
    file_list = os.listdir(path)
    dic_file_list = {}
    for file_name in file_list:
        num = int(file_name[:file_name.find('장')])
        dic_file_list[num] = file_name
    del file_list

    for song_num in worship_songs_copy:
        shutil.copy(path + dic_file_list[song_num], './output/' + '찬송가' + str(song_num).zfill(3) + '장.ppt')


# 예배에 쓸 교독문 복사 함수
def word_copy(bible_words_copy):
    path = "./교독문/"
    file_list = os.listdir(path)
    dic_file_list = {}
    for file_name in file_list:
        num = int(file_name[:file_name.find('번')])
        dic_file_list[num] = file_name
    del file_list

    for word in bible_words_copy:
        shutil.copy(path + dic_file_list[word], './output/' + '교독문' + str(word).zfill(3) + '번.ppt')


# 한글 성경구절 추출 함수
def word_extract_kor(find_words_extracted):
    # 성경 책이름과 약어 mapping
    pd_bookname = pd.read_csv('./bible_bookname.csv', encoding='euc-kr')
    dict_books = dict(zip(pd_bookname['약어'], pd_bookname['풀네임']))
    # dict_books_niv = dict(zip(pd_bookname['약어'], pd_bookname['Eng_fullname']))
    pd_bookname = pd_bookname.set_index('약어')

    # return할 성경구절 list
    words_list_result = []

    for find_word in find_words_extracted:
        if find_word[1].isdigit():
            book_name = find_word[0]  # 책 약어줄
            lines_num = find_word[1:]  # 구절 장:절
        else:
            book_name = find_word[0:2]  # 책 약어줄
            lines_num = find_word[2:]  # 구절 장:절

        file_name = './개역개정-text/' + dict_books[book_name] + '.txt'

        # 말씀의 주소 출력
        words_head = pd_bookname.loc[book_name]['book_name_Kor'] \
                     + " (" + pd_bookname.loc[book_name]['book_name_Eng'] + ") " + lines_num
        # print(words_head)

        words = []
        with open(file_name, 'r') as file:
            for text in file:
                start_index = text.find(' ', 1)
                words.append([text[0:start_index], text[start_index + 1:].strip('\n')])
                # if a == 10: break

        df_words = pd.DataFrame(words, columns=['lines', 'phrase'])
        del words
        df_words = df_words.set_index('lines')

        if '~' in find_word:  # 여러구절을 세트로
            loc0 = find_word.find(':')
            loc1 = find_word.find('~')
            # print(find_word[loc0+1:loc1], find_word[loc1+1:])
            start_num = int(find_word[loc0 + 1:loc1])
            end_num = int(find_word[loc1 + 1:])
            # print(start_num, end_num)
            for i in range(start_num, end_num + 1):
                find_word_multi = find_word[:loc0 + 1] + str(i)
                strings = df_words.loc[find_word_multi]['phrase']
                # words_out = str(i) + '\t' + strings
                words_out_list = [words_head, i, strings]
                words_list_result.append(words_out_list)
                # print(words_out)

        else:  # 한구절씩
            strings = df_words.loc[find_word]['phrase']
            loc0 = find_word.find(':')
            line_num = int(find_word[loc0 + 1:])
            words_out_list = [words_head, line_num, strings]
            words_list_result.append(words_out_list)

            # print(line_num, '\t', strings, sep='')
        # print("\n")
    return words_list_result


# 영문 성결구절 추출 함수
def word_extract_eng(find_words):
    # 성경 책이름과 약어 mapping
    pd_bookname = pd.read_csv('./bible_bookname.csv', encoding='euc-kr')
    dict_books = dict(zip(pd_bookname['약어'], pd_bookname['풀네임']))
    dict_books_niv = dict(zip(pd_bookname['약어'], pd_bookname['Eng_fullname']))
    pd_bookname = pd_bookname.set_index('약어')

    # return할 성경구절 list
    words_list_result = []

    for find_word in find_words:
        if (find_word[1].isdigit()):
            book_name = find_word[0]
        else:
            book_name = find_word[0:2]

        file_name = './NIV_English_Bible/' + dict_books_niv[book_name] + '.txt'
        # print(dict_books_niv[book_name])

        # 말씀의 주소 출력
        # print(find_word)

        words = []
        with open(file_name, 'r') as file:
            for text in file:
                start_index = text.find(' ', 1)
                words.append([text[0:start_index], text[start_index + 1:].strip('\n')])
                # if a == 10: break

        df_words = pd.DataFrame(words, columns=['lines', 'phrase'])
        del words
        df_words = df_words.set_index('lines')

        if (find_word[1].isdigit()):
            # book_name = find_word[0]
            book_name_eng = find_word[0]
            find_word = find_word[1:]
        else:
            book_name_eng = find_word[:2]
            find_word = find_word[2:]
            # book_name = find_word[0:2]

        if '~' in find_word:  # 여러구절을 세트로
            loc0 = find_word.find(':')
            loc1 = find_word.find('~')
            # print(find_word[loc0+1:loc1], find_word[loc1+1:])
            start_num = int(find_word[loc0 + 1:loc1])
            end_num = int(find_word[loc1 + 1:])
            # print(start_num, end_num)
            for i in range(start_num, end_num + 1):
                find_word_multi = find_word[:loc0 + 1] + str(i)
                strings = df_words.loc[find_word_multi]['phrase']
                # print(str(i), '\t', strings)
                words_out_list = [find_word, i, strings]
                words_list_result.append(words_out_list)
        else:  # 한구절씩
            strings = df_words.loc[find_word]['phrase']
            loc0 = find_word.find(':')
            line_num = int(find_word[loc0 + 1:])
            # print(line_num, '\t', strings, sep='')
            words_out_list = [find_word, line_num, strings]
            words_list_result.append(words_out_list)
        # print("\n")
    return words_list_result


##########input######
# 찬송가
worship_songs = []
# 교독문
bible_words = []
find_words = ['시73:25~26','롬8:6','잠4:23','빌4:7','마14:28~30','사40:31']





# -를 ~로 통일
find_words_temp = []
for words_item in find_words:
    words_item = words_item.replace('-', '~')
    find_words_temp.append(words_item)
find_words = find_words_temp
del find_words_temp
#####################

# 찬송가 및 PPT 카피해서 ./temp 폴더에 모아두기
song_copy(worship_songs)
word_copy(bible_words)

# 성경구절 추출해서 ./temp/test.csv에 저장하기
words_list_result_kor = word_extract_kor(find_words)
words_list_result_eng = word_extract_eng(find_words)

# Dataframe으로 전환
df_words_results = pd.DataFrame(words_list_result_kor, columns=['words_list', 'line_num', 'Words_kor'])
df_words_results_eng = pd.DataFrame(words_list_result_eng, columns=['words_list', 'line_num', 'Words_eng'])
df_words_results['Words_eng'] = df_words_results_eng['Words_eng']
del (df_words_results_eng)

# csv 파일로 저장.
# df_words_results.to_csv("./output/test.csv", encoding='utf-8-sig')

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from matplotlib import font_manager

prs = Presentation()

prs.slide_height = Cm(19.05)
prs.slide_width = Cm(33.867)

blank_slide_layout = prs.slide_layouts[6]  # black slide

left_k = Cm(0.05)
top_k = Cm(2.0)
width_k = Cm(33)
height_k = Cm(7)

left_e = Cm(0.05)
top_e = Cm(15.0)


for nline in df_words_results.index:
    # 슬라이드 추가 및 슬라이드 배경색 => 검정으로
    slide = prs.slides.add_slide(blank_slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # PPT 상단에 성경구절의 장:절 표시
    txBox_name = slide.shapes.add_textbox(Cm(0.05), Cm(0.05), Cm(24.0), Cm(1.5))  # 위치 (시작점x, 시작점y, 너비, 높이)
    tf_name = txBox_name.text_frame

    tf_name.text = df_words_results.loc[nline, 'words_list']  # 성경구절 주소 '고린도후서(2 Corinthians) 11장'
    print(tf_name.text)
    tf_name.paragraphs[0].font.size = Pt(32)
    tf_name.paragraphs[0].line_spacing = 1.0
    tf_name.paragraphs[0].font.name = font_manager.FontProperties(fname="c:/Windows/Fonts/H2HDRM.ttf").get_name()
    tf_name.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)  # 회색
    tf_name.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    # tf_name.fit_text()

    # 한글 구절
    txBox = slide.shapes.add_textbox(left_k, top_k, width_k, height_k)
    tf = txBox.text_frame

    #print용
    tf.text = str(df_words_results.loc[nline, 'line_num']) + "\t" + df_words_results.loc[nline, 'Words_kor']
    print(tf.text)
    
    
    tf.text = str(df_words_results.loc[nline, 'line_num']) + "\n" + df_words_results.loc[nline, 'Words_kor']
    tf.word_wrap = None

    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].line_spacing = 1.0
    tf.paragraphs[0].font.name = 'Times'
    tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)

    tf.paragraphs[1].font.size = Pt(32)
    #tf.paragraphs[1].line_spacing = Pt(40)
    tf.paragraphs[1].line_spacing = 1.0
    tf.paragraphs[1].font.name = font_manager.FontProperties(fname="c:/Windows/Fonts/H2HDRM.ttf").get_name()
    tf.paragraphs[1].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    tf.paragraphs[1].font.color.rgb = RGBColor(255, 255, 255)

    # 영문 구절
    txBox_e = slide.shapes.add_textbox(left_e, top_e, width_k, height_k)
    tf_e = txBox_e.text_frame
    
    # print용
    tf_e.text = str(df_words_results.loc[nline, 'line_num']) + "\t" + df_words_results.loc[nline, 'Words_eng']
    print(tf_e.text)
    
    tf_e.text = str(df_words_results.loc[nline, 'line_num']) + "\n" + df_words_results.loc[nline, 'Words_eng']
    tf_e.word_wrap = None

    tf_e.paragraphs[0].font.size = Pt(24)
    #tf_e.paragraphs[0].line_spacing = Pt(40)
    tf_e.paragraphs[0].line_spacing = 1.0
    tf_e.paragraphs[0].font.name = 'Times'
    tf_e.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    tf_e.paragraphs[0].font.color.rgb = RGBColor(255, 255, 0)

    tf_e.paragraphs[1].font.size = Pt(28)
    # tf_e.paragraphs[1].line_spacing = Pt(40)
    tf_e.paragraphs[0].line_spacing = 1.0
    tf_e.paragraphs[1].font.name = 'Times'
    tf_e.paragraphs[1].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    tf_e.paragraphs[1].font.color.rgb = RGBColor(255, 255, 255)

time_now = datetime.datetime.now()

hour_minute = time_now.hour * 100 + time_now.minute
filename = './output/주일말씀' + str(hour_minute) +'.pptx'

#prs.save('./output/말씀.pptx')
prs.save(filename)
for item in words_list_result_kor:
    print (item[1], item[2])
    
for item in words_list_result_eng:
    print (item[1], item[2])
